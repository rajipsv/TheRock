#!/usr/bin/env python3
"""Generate hackathon submission PowerPoint for Change Impact Agent."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parents[1] / "AGENTS_030_Hackathon_Deck.pptx"

CODE_URL = "https://github.com/rajipsv/TheRock/tree/feature/change-impact-agent/agents/change-impact-agent"
REPO_URL = "https://github.com/rajipsv/TheRock/tree/feature/change-impact-agent"
SAMPLES_URL = "https://github.com/rajipsv/TheRock/tree/feature/change-impact-agent/agents/change-impact-agent/sample-runs"
DEMO_VIDEO = "https://your-demo-video-url"  # replace after recording


def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def add_section_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], font_size: int = 18) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_size)


def add_table_slide(prs: Presentation, title: str, rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    cols, row_count = len(rows[0]), len(rows)
    left, top, width, height = Inches(0.5), Inches(1.4), Inches(9), Inches(0.35 * row_count)
    table = slide.shapes.add_table(row_count, cols, left, top, width, height).table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
            for paragraph in table.cell(r, c).text_frame.paragraphs:
                paragraph.font.size = Pt(14)


def main() -> None:
    prs = Presentation()

    # --- 1. Introduction ---
    add_section_slide(prs, "1. Introduction")

    add_title_slide(
        prs,
        "AGENTS_030 — Change Impact Agent for TheRock",
        "Evaluates infrastructure changes • predicts impact • suggests safer rollout",
    )

    add_bullet_slide(
        prs,
        "Hackathon use case (AGENTS_030)",
        [
            "Create an agent that evaluates planned infrastructure changes",
            "(patches, config updates, deployments), predicts potential impacts",
            "using historical incidents, and suggests safer rollout strategies.",
            "",
            "Our implementation: change-impact-agent (pre-merge) + log-analysis-agent (post-CI incidents).",
        ],
        font_size=17,
    )

    add_bullet_slide(
        prs,
        "Short description — What, Why, How",
        [
            "What: Agent evaluates planned ROCm/TheRock changes (patches, CI config, packaging) before merge.",
            "Why: Infrastructure PRs are high-risk — wrong rollout scope wastes GPU CI or ships regressions.",
            "How: Deterministic pipeline (manifest → paths → content → topology → rollout); optional vLLM brief.",
            "Output: report.json, report.html, executive_summary.md with severity, labels, rollout_strategy.",
        ],
    )

    add_table_slide(
        prs,
        "Team",
        [
            ["Field", "Value"],
            ["Team name", "TheRock Agents (AGENTS_030)"],
            ["Lead", "Rajeswari — your.email@example.com — Developer / architect"],
            ["Member 2", "[Name] — [email] — [role]"],
            ["Member 3", "[Name] — [email] — [role]"],
        ],
    )

    # --- 2. Problem & Context ---
    add_section_slide(prs, "2. Problem & Context")

    add_bullet_slide(
        prs,
        "Problem statement",
        [
            "Planned infrastructure changes (patches, config, deployments) land as TheRock PRs daily.",
            "Patches: submodule / superrepo SHA bumps (e.g. rocm-libraries → miopen, hipblaslt).",
            "Config updates: CI test-matrix timeouts, disabled jobs, artifact TOML, third-party CMake.",
            "Deployments: which test:* suites and test_filter:* depth to run across multi_arch_ci GPU farms.",
            "Teams lack a unified pre-merge view of blast radius, past failure patterns, and safe rollout order.",
        ],
        font_size=16,
    )

    add_bullet_slide(
        prs,
        "Current solutions, gap, difficulty",
        [
            "Existing: manifest-diff.yml (CI), assistant-librarian (bump PRs), manual review, ad-hoc log grep.",
            "Gap: No agent ties planned change → topology blast radius → historical CI incidents → rollout plan.",
            "Difficulty: Superrepo drill-down needs GitHub API; topology is domain-specific; incidents live in scattered logs.",
            "Importance: MI300-class CI is expensive — staged rollout (canary gfx) is mandatory at ROCm scale.",
        ],
        font_size=16,
    )

    add_bullet_slide(
        prs,
        "Target users & hackathon mapping",
        [
            "Users: ROCm/TheRock maintainers, CI owners, release/integration engineers, fork contributors.",
            "Stakeholders: AMD ROCm platform team, open-source reviewers, hackathon judges evaluating agent tooling.",
            "Market: Enterprise GPU software stacks with monorepo/superrepo CI (ROCm, large ML infra repos).",
            "Hackathon challenge: AGENTS_030 — AI agents for developer productivity on AMD MI300 / TheRock ecosystem.",
        ],
    )

    # --- 3. Solution Overview ---
    add_section_slide(prs, "3. Solution Overview")

    add_table_slide(
        prs,
        "Use case → implementation mapping",
        [
            ["Use case requirement", "How we deliver it", "Status"],
            [
                "Evaluate patches",
                "Manifest + superrepo component diff (--full-manifest)",
                "Done",
            ],
            [
                "Evaluate config updates",
                "content_diff.py — CI matrix timeouts, artifact TOML",
                "Done",
            ],
            [
                "Evaluate deployments",
                "ci_mapping.py — test:* labels, test_filter depth",
                "Done",
            ],
            [
                "Predict potential impacts",
                "impact_graph.py + BUILD_TOPOLOGY.toml severity score",
                "Done",
            ],
            [
                "Historical incidents",
                "log-analysis-agent failure KB (patterns + resolutions)",
                "Sibling agent",
            ],
            [
                "Safer rollout strategies",
                "rollout_strategy — canary gfx + staged test depth",
                "Done",
            ],
        ],
    )

    add_bullet_slide(
        prs,
        "Solution architecture / workflow",
        [
            "1. EVALUATE: PR or git range → manifest, paths, CI content, superrepo components.",
            "2. PREDICT: Map changes onto BUILD_TOPOLOGY.toml → severity, blast radius, build stages.",
            "3. INCIDENTS: Post-CI, log-analysis-agent matches failure_kb patterns; feeds future triage.",
            "4. ROLLOUT: rollout_strategy + test:* / test_filter:* labels (canary gfx family first).",
            "5. REPORT: report.json + report.html + executive_summary.md (optional vLLM brief).",
        ],
        font_size=16,
    )

    add_bullet_slide(
        prs,
        "AI approach (agents, not chat-only)",
        [
            "Agent-assisted pipeline: fixed tool order for reliability; LLM optional for prose only.",
            "Pre-merge agent (change-impact): evaluate → predict topology impact → recommend rollout.",
            "Post-failure agent (log-analysis): vectorless/keyword KB retrieval over historical CI incidents.",
            "Optional vLLM + Qwen3-30B-A3B: reviewer brief from pre-computed JSON (summarize.py).",
            "No fine-tuning; no vector RAG for topology — BUILD_TOPOLOGY.toml is source of truth.",
        ],
        font_size=16,
    )

    add_bullet_slide(
        prs,
        "Safer rollout strategies (built-in)",
        [
            "CRITICAL/HIGH severity → canary one GPU family (gfx110X) then full multi_arch_ci matrix.",
            "MEDIUM → canary gfx + component-specific test:* labels (e.g. test:miopen).",
            "CI-only / packaging → test_filter:quick + manifest-diff sibling job.",
            "HIGH+ may add ci:run-all-archs label recommendation.",
            "Advisory only — humans apply labels; agent does not trigger CI or deploy.",
        ],
        font_size=16,
    )

    add_bullet_slide(
        prs,
        "Key technologies",
        [
            "Python 3.10+, Git, GitHub REST API, Jinja2 HTML reports.",
            "vLLM on AMD MI300X (Qwen/Qwen3-30B-A3B), OpenAI client, optional Ollama.",
            "GitHub Actions: change-impact-agent.yml, change-impact-upstream-scan.yml.",
            "pytest unit tests (15+ for change-impact); Jupyter hackathon_demo.ipynb.",
        ],
    )

    add_bullet_slide(
        prs,
        "Datasets & historical incidents",
        [
            "Infrastructure source of truth: ROCm/TheRock — BUILD_TOPOLOGY.toml, manifests, CI scripts.",
            "Validated change sets: PR #5572 (config), #5688 (config+packaging), #5480 (packaging), #5718 (patch).",
            "Historical incidents: log-analysis-agent/knowledge/patterns.json + resolutions.jsonl.",
            "Real GHA log corpus: upstream run-27697860238 (kfdtest PR #8864); keyword + optional FAISS hybrid retrieval.",
            "Committed sample-runs/ — full reports per PR for judges (no live API needed).",
        ],
        font_size=16,
    )

    add_bullet_slide(
        prs,
        "Built during hackathon",
        [
            "change-impact-agent: analyze.py, summarize.py, upstream_pr_scan.py, topology_audit.py, 7 analysis bridges.",
            "log-analysis-agent: analyze_log.py, presets, GitHub log fetch, GHA workflows.",
            "document-comparison-agent & earnings-ir-agent: additional AMD TCS demos ported to agents/.",
            "Docs: README, SCOPE, PRESENTATION, notebook walkthrough, pre-generated sample reports.",
        ],
        font_size=16,
    )

    # --- 4. Details ---
    add_section_slide(prs, "4. Technical Details")

    add_bullet_slide(
        prs,
        "Models",
        [
            "Primary (optional LLM): Qwen/Qwen3-30B-A3B via vLLM — served-model-name Qwen3-30B-A3B.",
            "Inference: OpenAI-compatible /v1/chat/completions; temperature 0.2 for summaries.",
            "Fine-tuning: None — not required; core value is deterministic analysis.",
            "Default path: Template executive summary (zero GPU, instant).",
        ],
    )

    add_bullet_slide(
        prs,
        "Performance & pipeline",
        [
            "analyze.py (typical PR, no --full-manifest): ~30–90 s (git + GitHub API).",
            "--full-manifest superrepo drill-down (#5718): several minutes (GitHub compare + commits).",
            "summarize.py template: <1 s; vLLM brief: ~30–120 s depending on report size.",
            "Concurrency: CLI/batch via upstream_pr_scan.py; GHA dispatch for upstream scan.",
            "End-to-end: PR number → HTML report in one command.",
        ],
        font_size=16,
    )

    add_bullet_slide(
        prs,
        "GPU usage (MI300 / vLLM)",
        [
            "GPU required only for optional LLM summaries and notebook demos — not for core analyze.py.",
            "vLLM launch: Qwen3-30B-A3B, port 8000, --enable-auto-tool-choice --tool-call-parser hermes.",
            "Memory: ~30B MoE model — fits MI300 class GPU with vLLM; watch rocm-smi during serve.",
            "Energy: Targeted test:* recommendations reduce wasted multi_arch_ci GPU hours (indirect savings).",
        ],
    )

    add_table_slide(
        prs,
        "Use case evidence — real upstream PRs",
        [
            ["PR", "Change type", "Predicted impact", "Safer rollout"],
            [
                "5572",
                "Config (CI timeout)",
                "MEDIUM — miopen matrix",
                "Canary gfx + test:miopen, test_filter:quick",
            ],
            [
                "5688",
                "Config + packaging",
                "MEDIUM — hipDNN path",
                "Canary gfx + test:hipdnn, test_filter:quick",
            ],
            ["5480", "Packaging / third-party", "Low packaging risk", "test_filter:quick"],
            [
                "5718",
                "Patch (superrepo bump)",
                "Multi-component blast radius",
                "Canary gfx + component-scoped test:*",
            ],
        ],
    )

    # --- 5. Summary ---
    add_section_slide(prs, "5. Summary")

    add_bullet_slide(
        prs,
        "Expected impact & value (use case outcomes)",
        [
            "Evaluate faster: patches, config, and deployment scope auto-classified per PR.",
            "Predict smarter: topology blast radius + severity score before any GPU CI runs.",
            "Learn from incidents: log-analysis-agent KB closes the loop when CI fails.",
            "Roll out safer: canary gfx + staged test_filter (quick → standard → full).",
            "Productivity: minutes of triage → seconds; fewer wasted MI300 CI hours.",
        ],
        font_size=16,
    )

    add_bullet_slide(
        prs,
        "Innovation / key differentiators",
        [
            "End-to-end infrastructure change agent — not just diff, but rollout + CI deployment plan.",
            "Three change classes in one pipeline: patches, config updates, packaging/deployments.",
            "Topology-aware prediction (BUILD_TOPOLOGY.toml) + historical incident KB (sibling agent).",
            "Deterministic-first: authoritative JSON; LLM only rephrases, never invents facts.",
            "Validated on live ROCm/TheRock upstream PRs with committed sample-runs/.",
        ],
        font_size=16,
    )

    add_bullet_slide(
        prs,
        "Future work",
        [
            "Wire failure_kb into pre-merge analyze.py (incident-aware rollout warnings).",
            "Merge agent workflows upstream to ROCm/TheRock; scheduled upstream PR scan.",
            "Auto-append resolutions.jsonl from post-CI triage → richer historical prediction.",
            "Deeper multi_arch_ci canary integration; maintainer-gated label application.",
        ],
        font_size=16,
    )

    add_bullet_slide(
        prs,
        "Links",
        [
            f"Code (agent): {CODE_URL}",
            f"Repo branch: {REPO_URL}",
            f"Sample reports: {SAMPLES_URL}",
            f"Demo video: {DEMO_VIDEO}  ← replace after recording",
            "Notebook: agents/change-impact-agent/notebook/hackathon_demo.ipynb",
        ],
        font_size=14,
    )

    add_title_slide(prs, "Thank you", "Questions?")

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
